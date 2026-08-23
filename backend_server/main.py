import os
import io
import json
import hashlib
import base64
import mimetypes
from datetime import datetime, timedelta, timezone
from enum import Enum
from typing import List, Optional, Dict,Literal, Any,Union
from fastapi import FastAPI, Depends, HTTPException, status, Response, Request, UploadFile, File, Form, Header, Cookie, Query 
from fastapi.middleware.cors import CORSMiddleware 
from fastapi.security import OAuth2PasswordBearer
from fastapi.responses import JSONResponse
from pydantic import BaseModel, EmailStr,Field
from langchain_text_splitters import RecursiveCharacterTextSplitter
from passlib.context import CryptContext
from jose import JWTError, jwt
import supabase
from supabase import create_client, Client
import groq 
import uvicorn
import redis
import requests
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from pdf2image import convert_from_bytes
import pytesseract
import httpx
import uuid
import psycopg2
from psycopg2.extras import RealDictCursor 
from dotenv import load_dotenv

load_dotenv()

SECRET_KEY = os.getenv("JWT_SECRET_KEY", "SUPER_SECRET_JWT_KEY_CHANGE_IN_PRODUCTION")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

LOCAL_SESSION_TTL = 3600   
SUPABASE_URL = os.getenv("SUPABASE_URL", "https://your-supabase-project.supabase.co")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "your-supabase-service-role-key")

NEON_DATABASE_URL = os.getenv("NEON_DATABASE_URL", "postgresql://user:password@ep-host.region.aws.neon.tech/dbname?sslmode=require")

try:
    supabase_client: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
except Exception as e:
    print(f"Warning: Supabase initialization failed: {e}")
    supabase_client = None

# Support REDIS_URL (e.g. redis://redis_db:6379 set by docker-compose) or
# individual REDIS_HOST / REDIS_PORT for local development.
_redis_url = os.getenv("REDIS_URL")
if _redis_url:
    # Parse redis://host:port
    _parts = _redis_url.replace("redis://", "").split(":")
    REDIS_HOST = _parts[0]
    REDIS_PORT = int(_parts[1]) if len(_parts) > 1 else 6379
else:
    REDIS_HOST = os.getenv("REDIS_HOST", "localhost")
    REDIS_PORT = int(os.getenv("REDIS_PORT", 6379))
try:
    redis_client = redis.Redis(host=REDIS_HOST, port=REDIS_PORT, db=0, decode_responses=True)
    redis_client.ping()
except Exception as e:
    print(f"Warning: Redis connection failed: {e}")
    redis_client = None

GROQ_API_KEY = os.getenv("GROQ_API_KEY", "your-groq-api-key")
JINA_API_KEY = os.getenv("JINA_API_KEY", "your-jina-api-key")

GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID", "your-google-client-id")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET", "your-google-client-secret")

def get_neon_connection():
    try:
        return psycopg2.connect(NEON_DATABASE_URL)
    except Exception as e:
        print(f"Warning: Neon DB connection failed: {e}")
        return None
        
groq_client = groq.Groq(api_key=GROQ_API_KEY)

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

Oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/v1/auth/login", auto_error=False)

app = FastAPI(
    title="Global RAG Application API",
    version="2.1.0",
    description="Enterprise RAG backend with LangChain chunking, Supabase Delta Registries, and RBAC."
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173","http://13.201.77.82","http://13.201.77.82:80"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class Role(str, Enum):
    SUPER_ADMIN = "Super Admin"
    ADMIN = "Admin"
    USER = "User"

class DocumentType(str, Enum):
    PLAIN_TEXT = "plain text"
    PDF_TEXT = "pdf text"
    PDF_SCANNED = "pdf scanned"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    IMAGE = "image"
    HTML = "html"
    UNKNOWN = "unknown"
    
class UploadMode(str, Enum):
    GLOBAL = "global"
    LOCAL = "local"    

class OrgSettingsUpdate(BaseModel):
    allow_user_global_upload: bool

class OrgSettingsResponse(BaseModel):
    org_id: str
    allow_user_global_upload: bool

class UserSignup(BaseModel):
    email: EmailStr
    password: str = Field(..., min_length=8)
    org_id: str
    tenant_id: Optional[str] = None
    Role : str 

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class TokenClaims(BaseModel):
    user_id: str
    email: str
    role: Role
    org_id: Optional[str] = None
    tenant_id: Optional[str] = None
    exp: Optional[int] = None
    iat: Optional[int] = None
    
class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    role: Role
    org_id: Optional[str] = None

class DocumentMetadata(BaseModel):
    file_name: str
    file_size: int
    file_hash: str
    total_pages: int
    upload_type: UploadMode
    uploaded_by: str
    uploaded_at: str
    org_id: str
    status: str = "processing"

class ChunkMetadata(BaseModel):
    chunk_id: str
    document_id: str
    page_number: int
    chunk_index: int
    org_id: str
    upload_type: UploadMode
    text_preview: str
    embedding_model: str        

class QueryMode(str, Enum):
    GLOBAL = "global"   
    LOCAL  = "local"    
    BOTH   = "both"     

class RAGRequest(BaseModel):
    query: str
    session_id: Optional[str] = None  
    upload_mode: QueryMode = QueryMode.GLOBAL
    top_k: int = 5
    vector_weight: float = 0.7     
    keyword_weight: float = 0.3
    language: str = "English"    
    system_prompt: Optional[str] = None

class SourceResult(BaseModel):
    chunk_id: str
    document_id: str 
    document_name: str
    page_number: int
    chunk_index: int
    similarity_score: float
    text_preview: str            
    org_id: str
    upload_mode: str
    document_url: str

class RAGResponse(BaseModel):
    answer: str
    query: str
    language: str
    query_mode: QueryMode
    sources: List[SourceResult]
    total_sources_found: int
    generated_by: str
    org_id: str
    queried_at: str
    session_id: str           

    
class UserInfo(BaseModel):
    user_id: str
    email: str
    role: Role
    org_id: Optional[str] = None
    created_at: str
        
def hash_password(password: str) -> str:
    return pwd_context.hash(password)

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return pwd_context.verify(plain_password, hashed_password)

def create_jwt_token(claims: Dict[str, Any], expires_delta: timedelta) -> str:
    to_encode = claims.copy()
    expire = datetime.now(timezone.utc) + expires_delta
    to_encode.update({
        "exp": int(expire.timestamp()),
        "iat": int(datetime.now(timezone.utc).timestamp())
    })
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def decode_jwt_token(token: str) -> Dict[str, Any]:
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        return payload
    except JWTError:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Token invalid or expired"
        )
        
def file_hash(file_bytes: bytes) -> str:
    return hashlib.sha256(file_bytes).hexdigest()
    
def page_hash(page_content: Union[str, bytes]) -> str:
    if isinstance(page_content, str):
        page_content = page_content.encode('utf-8')
    return hashlib.md5(page_content).hexdigest()
    
def generate_chunk_id() -> str:
    return f"chunk_{uuid.uuid4().hex[:12]}"

def generate_document_id() -> str:
    return f"doc_{uuid.uuid4().hex[:12]}"                                                                            
def detect_document_type(filename: str, file_bytes: bytes) -> DocumentType:
    filename_lower = filename.lower()
    if filename_lower.endswith('.txt') or filename_lower.endswith('.md'):
        return DocumentType.PLAIN_TEXT
    elif filename_lower.endswith('.pdf'):
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            first_page = reader.pages[0]
            text = first_page.extract_text()
            if not text or len(text.strip()) < 50:
                return DocumentType.PDF_SCANNED
            return DocumentType.PDF_TEXT
        except:
            return DocumentType.PDF_SCANNED
    elif filename_lower.endswith('.docx'):
        return DocumentType.DOCX
    elif filename_lower.endswith(('.xlsx', '.csv')):
        return DocumentType.XLSX
    elif filename_lower.endswith('.pptx'):
        return DocumentType.PPTX
    elif filename_lower.endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp')):
        return DocumentType.IMAGE
    elif filename_lower.endswith('.html') or filename_lower.endswith('.htm'):
        return DocumentType.HTML
    return DocumentType.UNKNOWN
    
def sentence_chunking(text: str) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n"]
    )
    return splitter.split_text(text)
    
def recursive_chunking(text: str) -> List[str]:
    recursive_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    length_function=len,
    separators=["\n\n", "\n", " ", ""]
    )
    chunks = recursive_splitter.split_text(text)
    return chunks 

def check_org_global_upload_setting(org_id: str) -> bool:
    if not supabase_client:
        return False
    try:
        response = (
            supabase_client.table("organization_settings")
            .select("allow_user_global_upload")
            .eq("org_id", org_id)
            .execute()
        )
        if response.data:
            return response.data[0].get("allow_user_global_upload", False)
    except Exception as e:
        print(f"Error checking org settings: {e}")
    return False


def docx_chunking(text: str, max_chunk_size: int = 1500) -> List[str]:
    chunks = []
    current_section = ""
    for line in text.split('\n'):
        if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
            if current_section:
                chunks.append(current_section.strip())
            current_section = line
        else:
            current_section += "\n" + line
            if len(current_section) > max_chunk_size:
                chunks.append(current_section.strip())
                current_section = "" 
    if current_section:
        chunks.append(current_section.strip())
    return [c for c in chunks if len(c) > 30]

def row_chunking(text: str, rows_per_chunk: int = 20) -> List[str]:
    if not text.strip():
        return []
    lines = text.split('\n')
    chunks = []
    current_chunk = []
    for line in lines:
        if not line.strip():
            continue
        current_chunk.append(line)
        if len(current_chunk) == rows_per_chunk:
            chunks.append("\n".join(current_chunk))
            current_chunk = []
    if current_chunk:
        chunks.append("\n".join(current_chunk))
    return chunks
    
def ppt_chunking(text: str) -> List[str]:
    slides = text.split('---SLIDE_BREAK---')
    chunk_of_slide = [s.strip() for s in slides if len(s.strip()) > 30]
    return chunk_of_slide
        
def chunk_tag_aware(text: str, max_chunk_size: int = 1500) -> List[str]:
    chunks = []
    current_section = ""
    lines = [line for line in text.split('\n') if line.strip()]
    for line in lines:
        if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
            if current_section:
                chunks.append(current_section.strip())
            current_section = line
        else:
            current_section += "\n" + line
            if len(current_section) > max_chunk_size:
                chunks.append(current_section.strip())
                current_section = "" 
    if current_section:
        chunks.append(current_section.strip())
    return [c for c in chunks if len(c) > 30]
        
def extract_pdf_text(file_bytes: bytes) -> tuple[str, int]:
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        page_count = len(reader.pages)
        all_pages = []
        for page_number, page in enumerate(reader.pages):
            extracted_text = page.extract_text()
            if extracted_text and len(extracted_text.strip()) > 50:
                page_text = extracted_text
            else:
                try:
                    images = convert_from_bytes(
                        file_bytes,
                        dpi=300,
                        first_page=page_number + 1,
                        last_page=page_number + 1
                    )
                    if images:
                        page_text = pytesseract.image_to_string(
                            images[0],
                            lang="eng",
                            config="--oem 3 --psm 6"
                        )
                    else:
                        page_text = ""
                except Exception as ocr_err:
                    page_text = extracted_text
            all_pages.append(page_text.strip())
        final_text = "\n---PAGE_BREAK---\n".join(all_pages)
        return final_text, page_count
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF extraction failed: {str(e)}")

def extract_docx_text(file_bytes: bytes) -> tuple[str, int]:
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        text_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name.lower()
            if style_name.startswith('heading 1'):
                text = f"# {text}"
            elif style_name.startswith('heading 2'):
                text = f"## {text}"
            elif style_name.startswith('heading 3'):
                text = f"### {text}"
            text_parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                if any(row_data):
                    text_parts.append(" | ".join(row_data))
        full_text = "\n".join(text_parts)
        block_count = len(doc.paragraphs) + sum(len(t.rows) for t in doc.tables)
        return full_text, block_count
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX extraction failed: {str(e)}") 
               
def extract_xlsx_text(file_bytes: bytes) -> tuple[str, int]:
    full_text_lines = []
    total_rows = 0
    try:
        try:
            excel_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
        except Exception:
            excel_dict = {"Sheet1": pd.read_csv(io.BytesIO(file_bytes))}
        for sheet_name, df in excel_dict.items():
            df = df.fillna("")
            total_rows += len(df)
            columns = df.columns.tolist()
            for index, row in df.iterrows():
                row_parts = [f"Sheet: {sheet_name}"]
                for col in columns:
                    val = str(row[col]).strip()
                    if val:   
                        row_parts.append(f"{col}: {val}")
                full_text_lines.append(" | ".join(row_parts))
        return "\n".join(full_text_lines), total_rows
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Spreadsheet extraction failed: {str(e)}")

def extract_text_from_shape(shape) -> str:
    text = ""
    if hasattr(shape, "text") and shape.text:
        text += shape.text + "\n"
    if shape.has_table:
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.replace('\n', ' ').strip()
                if cell_text:
                    row_data.append(cell_text)
            text += " | ".join(row_data) + "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            text += extract_text_from_shape(child_shape)
    return text

def extract_pptx_text(file_bytes: bytes) -> tuple[str, int]:
    try:
        prs = PptxPresentation(io.BytesIO(file_bytes))
        full_text = ""
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                slide_text += extract_text_from_shape(shape)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text += f"\nSpeaker Notes:\n{notes}\n"
            full_text += slide_text + "\n---SLIDE_BREAK---\n"
        return full_text, len(prs.slides)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX extraction failed: {str(e)}")            

def extract_html_text(file_bytes: bytes) -> tuple[str, int]:
    try:
        html_text = file_bytes.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html_text, 'html.parser')
        for tag in soup.find_all(['nav', 'footer', 'script', 'style', 'aside', 'header', 'iframe']):
            tag.decompose()
        for level in range(1, 7):
            for header in soup.find_all(f'h{level}'):
                markdown_header = f"\n{'#' * level} {header.get_text(strip=True)}\n"
                header.string = markdown_header
        for table in soup.find_all('table'):
            table_lines = []
            for row in table.find_all('tr'):
                cells = row.find_all(['td', 'th'])
                row_data = [cell.get_text(strip=True).replace('\n', ' ') for cell in cells]
                if any(row_data):
                    table_lines.append(" | ".join(row_data))
            if table_lines:
                table.insert_after("\n" + "\n".join(table_lines) + "\n")
                table.decompose()
        text = soup.get_text(separator='\n', strip=True)
        return text, 1
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"HTML extraction failed: {str(e)}")
        
def extract_plainfile_text(file_bytes : bytes):
    text = file_bytes.decode('utf-8')
    return text, 1
        
ALLOWED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp"}

IMAGE_EXTENSIONS = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}

def get_image_mime_type(filename: str) -> str:
    ext = os.path.splitext(filename.lower())[1]
    return IMAGE_EXTENSIONS.get(ext, mimetypes.guess_type(filename)[0] or "image/jpeg")

def read_image_to_base64(file_bytes: bytes, content_type: str) -> tuple[str, int]:
    if content_type not in ALLOWED_IMAGE_TYPES:
        raise HTTPException(status_code=400, detail=f"Unsupported image type: {content_type}. Allowed: {ALLOWED_IMAGE_TYPES}")
    try:
        base64_encoded = base64.b64encode(file_bytes).decode("utf-8")
        data_url = f"data:{content_type};base64,{base64_encoded}"
        return data_url, 1
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Image encoding failed: {str(e)}")    

def extract_text_by_type(doc_type: DocumentType, file_bytes: bytes, filename: str) -> tuple[str, int]:
    if doc_type == DocumentType.PLAIN_TEXT:
        return extract_plainfile_text(file_bytes)
    elif doc_type in [DocumentType.PDF_TEXT, DocumentType.PDF_SCANNED]:
        return extract_pdf_text(file_bytes)
    elif doc_type == DocumentType.DOCX:
        return extract_docx_text(file_bytes)
    elif doc_type == DocumentType.XLSX:
        return extract_xlsx_text(file_bytes)
    elif doc_type == DocumentType.PPTX:
        return extract_pptx_text(file_bytes)
    elif doc_type == DocumentType.HTML:
        return extract_html_text(file_bytes)
    elif doc_type == DocumentType.IMAGE:
        mime_type = get_image_mime_type(filename)
        return read_image_to_base64(file_bytes, mime_type)
    else:
        text = file_bytes.decode('utf-8', errors='ignore')
        return text, 1

def chunk_by_type(doc_type: DocumentType, text: str) -> List[str]:
    if doc_type == DocumentType.IMAGE:
        return [text]
    elif doc_type == DocumentType.PLAIN_TEXT:
        return sentence_chunking(text)
    elif doc_type in [DocumentType.PDF_TEXT, DocumentType.PDF_SCANNED]:
        return recursive_chunking(text)
    elif doc_type == DocumentType.DOCX:
        return docx_chunking(text)
    elif doc_type == DocumentType.XLSX:
        return row_chunking(text, rows_per_chunk=20)
    elif doc_type == DocumentType.PPTX:
        return ppt_chunking(text)
    elif doc_type == DocumentType.HTML:
        return chunk_tag_aware(text)
    else:
        return recursive_chunking(text)                  

def get_jina_embeddings(texts: List[str], is_image: bool = False) -> List[List[float]]:
    url = "https://api.jina.ai/v1/embeddings"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {JINA_API_KEY}"
    }
    if is_image:
        input_data = [{"image": text} for text in texts]
    else:
        input_data = [{"text": text} for text in texts]   
    payload = {
        "model": "jina-clip-v2",
        "normalized": True,
        "dimension":1024,
        "input": input_data
    }
    try:
        response = requests.post(url, json=payload, headers=headers, timeout=20)
        if response.status_code == 200:
            res_json = response.json()
            return [item["embedding"] for item in res_json["data"]]
        else:
            return [[0.01 * (i + 1) for i in range(1024)] for _ in texts]
    except Exception:
        return [[0.01 * (i + 1) for i in range(1024)] for _ in texts]

def check_file_in_registry(file_hash_val: str, org_id: str) -> Optional[Dict[str, Any]]:
    if not supabase_client:
        return None
    try:
        response = (
            supabase_client.table("document_registry")
            .select("*")
            .eq("file_hash", file_hash_val)
            .eq("org_id", org_id)
            .execute()
        )
        if response.data and len(response.data) > 0:
            return response.data[0]
    except Exception as e:
        return f"Error checking document registry in Supabase: {e}"
    return None

def store_file_in_registry(
    file_hash_val: str, 
    org_id: str, 
    doc_id: str, 
    filename: str, 
    total_pages: int, 
    upload_type: str, 
    uploaded_by: str
) -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("document_registry").insert({
            "id": doc_id,
            "org_id": org_id,
            "file_name": filename,
            "file_hash": file_hash_val,
            "uploaded_by": uploaded_by,
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
            "total_pages": total_pages,
            "status": "processed"
        }).execute()
    except Exception as e:
        print(f"Error storing file in Supabase document registry: {e}")

def check_page_hash_in_registry(doc_id: str, page_num: int, page_hash_val: str) -> bool:
    if not supabase_client:
        return False
    try:
        response = (
            supabase_client.table("page_registry")
            .select("id")
            .eq("document_id", doc_id)
            .eq("page_number", page_num)
            .eq("page_hash", page_hash_val)
            .execute()
        )
        return len(response.data) > 0 if response.data else False
    except Exception as e:
        print(f"Error checking page registry in Supabase: {e}")
        return False

def store_page_hash(doc_id: str, page_num: int, page_hash_val: str, chunk_ids: List[str], org_id:str = "") -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("page_registry").insert({
            "document_id": doc_id,
            "org_id": org_id,
            "page_number": page_num,
            "page_hash": page_hash_val,
            "chunk_ids": chunk_ids,
            "is_indexed": True,
            "indexed_at": datetime.now(timezone.utc).isoformat()
        }).execute()
    except Exception as e:
        print(f"Error storing page hash in Supabase page registry: {e}")

def extract_raw_jwt(
    authorization: Optional[str] = Header(None),
    access_token_cookie: Optional[str] = Cookie(None, alias="access_token")
) -> str:
    if authorization:
        parts = authorization.split()
        if len(parts) == 2 and parts[0].lower() == "bearer":
            return parts[1]
        return authorization
    if access_token_cookie:
        return access_token_cookie
    raise HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Authentication JWT token missing."
    )

async def get_current_user(
    token: Optional[str] = Depends(extract_raw_jwt)
) -> TokenClaims:
    if not token:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="No token provided"
        )
    payload = decode_jwt_token(token)
    return TokenClaims(
        user_id=payload.get("user_id"),
        email=payload.get("email"),
        role=Role(payload.get("role", "USER")),
        org_id=payload.get("org_id"),
        tenant_id=payload.get("tenant_id"),
        exp=payload.get("exp")
    )

class RoleChecker:
    def __init__(self, allowed_roles: List[Role]):
        self.allowed_roles = allowed_roles
    def __call__(self, current_user: TokenClaims = Depends(get_current_user)) -> TokenClaims:
        if current_user.role not in self.allowed_roles:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail=f"Access denied. Required role: {[r.value for r in self.allowed_roles]}"
            )
        return current_user
        
def enforce_tenant_access(requested_org_id: str, current_user: TokenClaims):
    if current_user.role == Role.SUPER_ADMIN:
        return True  
    if current_user.org_id != requested_org_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Tenant Isolation Policy: Access to requested organisation data is forbidden"
        )
    return True
    
def decode_base64_to_image(data_url: str):
    if not data_url.startswith("data:") or ";base64," not in data_url:
        raise HTTPException(status_code=500, detail="Stored image format is invalid or corrupted")
    header, encoded = data_url.split(";base64,", 1)
    media_type = header.replace("data:", "", 1)
    try:
        image_bytes = base64.b64decode(encoded)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Could not decode stored image: {e}")
    return media_type, image_bytes

def store_image_global_neon(
    doc_id: str,
    org_id: str,
    filename: str,
    data_url: str,
    embedding: List[float],
    upload_mode: UploadMode,
    uploaded_by: str,
    role: str 
) -> None:
    conn = get_neon_connection()
    if not conn:
        raise HTTPException(status_code=500, detail="Neon DB connection failed")
    try:
        with conn.cursor() as cur:

            cur.execute("SELECT set_config('app.current_org_id', %s, true)", (org_id,))
            cur.execute("SELECT set_config('app.current_user_id', %s, true)", (uploaded_by,))
            cur.execute("SELECT set_config('app.current_role', %s, true)", (role,))
            
            cur.execute(
                """
                INSERT INTO image_store (
                    id, org_id, file_name, data_url,
                    embedding, upload_mode, uploaded_by, uploaded_at
                )
                VALUES (%s, %s, %s, %s, %s::vector, %s, %s, %s)
                ON CONFLICT (id) DO NOTHING
                """,
                (
                    doc_id, org_id, filename, data_url,
                    embedding, upload_mode.value, uploaded_by,
                    datetime.now(timezone.utc).isoformat()
                )
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        raise HTTPException(status_code=500, detail=f"Neon image store failed: {e}")
    finally:
        conn.close()

def store_image_local_redis(
    doc_id: str,
    org_id: str,
    filename: str,
    data_url: str,
    embedding: List[float],
    ttl_seconds: int = 3600       
) -> None:
    if not redis_client:
        raise HTTPException(status_code=500, detail="Redis connection unavailable")
    try:
        payload = json.dumps({
            "doc_id": doc_id, "org_id": org_id, "file_name": filename,
            "data_url": data_url, "embedding": embedding,
            "stored_at": datetime.now(timezone.utc).isoformat()
        })
        redis_client.setex(
            f"image:{org_id}:{doc_id}",  
            ttl_seconds,
            payload
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Redis image store failed: {e}")
        
def store_image_by_mode(
    doc_id: str,
    org_id: str,
    filename: str,
    data_url: str,
    embedding: List[float],
    upload_mode: UploadMode,
    uploaded_by: str,
    role: str # <-- Added param
) -> None:
    if upload_mode == UploadMode.GLOBAL:
        store_image_global_neon(
            doc_id, org_id, filename,
            data_url, embedding,
            upload_mode, uploaded_by, role
        )
    elif upload_mode == UploadMode.LOCAL:
        store_image_local_redis(
            doc_id, org_id, filename,
            data_url, embedding
        )
    else:
        raise HTTPException(status_code=400, detail=f"Unknown upload mode: {upload_mode}")                

def store_raw_file_local_redis(user_id: str, doc_id: str, filename: str, file_bytes: bytes) -> None:
    if not redis_client:
        return
    try:
        encoded_data = base64.b64encode(file_bytes).decode('utf-8')
        payload = json.dumps({
            "filename": filename,
            "data": encoded_data
        })
        redis_client.setex(f"local_raw:{user_id}:{doc_id}", LOCAL_SESSION_TTL, payload)
    except Exception as e:
        print(f"Failed to store raw file in Redis: {e}")

def store_raw_file_global_supabase(org_id: str, doc_id: str, filename: str, file_bytes: bytes) -> None:
    if not supabase_client:
        return
    try:
        path = f"{org_id}/{doc_id}/{filename}"
        supabase_client.storage.from_("global_documents").upload(
            file=file_bytes,
            path=path,
            file_options={"x-upsert": "true", "content-type": "application/octet-stream"}
        )
    except Exception as e:
        print(f"Failed to store raw file in Supabase: {e}")

class UploadConsent(BaseModel):
    upload_mode: UploadMode
    confirmed: bool = False         

class ConsentMessage(BaseModel):
    upload_mode: UploadMode
    title: str
    message: str
    confirm_label: str = "Got it"
    warning_label: Optional[str] = None

class UploadReport(BaseModel):
    doc_id: str
    file_name: str
    upload_mode: UploadMode
    org_id: str
    doc_type: str
    total_pages: int
    pages_newly_indexed: int       
    pages_skipped: int               
    chunks_created: int
    status: str
    uploaded_at: str                        

def run_delta_management(
    doc_id: str,
    org_id: str,
    filename: str,
    file_hash_val: str,
    raw_text: str,
    doc_type: DocumentType,
    chunks: List[str],
    uploaded_by: str,
    upload_mode: UploadMode,
    role: str 
) -> dict:
    if "---PAGE_BREAK---" in raw_text:
        pages = raw_text.split("---PAGE_BREAK---")
        pages = [p.strip() for p in pages if p.strip()]
    else:
        pages = [raw_text.strip()] if raw_text.strip() else []
    total_pages = len(pages) if pages else 1

    existing_doc = check_file_in_registry(file_hash_val, org_id)

    if existing_doc:
        registered_doc_id = existing_doc["id"]

        # ── Alias tracking: log if a different filename was used for the same content ──
        original_filename = existing_doc.get("file_name", "unknown")
        if filename != original_filename and supabase_client:
            try:
                supabase_client.table("audit_log").insert({
                    "event_type":        "duplicate_file_detected",
                    "doc_id":            registered_doc_id,
                    "alias_filename":    filename,           # e.g. Brd_v2.pdf
                    "original_filename": original_filename,  # e.g. Brd_v1.pdf
                    "file_hash":         file_hash_val,
                    "user_id":           uploaded_by,
                    "org_id":            org_id,
                    "timestamp":         datetime.now(timezone.utc).isoformat()
                }).execute()
            except Exception as e:
                print(f"Alias audit log write failed: {e}")

        pages_skipped = 0
        pages_newly_indexed = 0
        chunks_created = 0
        new_chunk_ids: List[str] = []

        for page_num, page_text in enumerate(pages, start=1):
            p_hash = page_hash(page_text)
            already_indexed = check_page_hash_in_registry(
                registered_doc_id, page_num, p_hash
            )
            if already_indexed:
                pages_skipped += 1
                continue

            is_image = doc_type == DocumentType.IMAGE
            page_chunks = [page_text] if is_image else chunk_by_type(doc_type, page_text)
            embeddings = get_jina_embeddings(page_chunks, is_image=is_image)

            conn = get_neon_connection()
            if conn:
                try:
                    with conn.cursor() as cur:
                        
                        cur.execute("SELECT set_config('app.current_org_id', %s, true)", (org_id,))
                        cur.execute("SELECT set_config('app.current_user_id', %s, true)", (uploaded_by,))
                        cur.execute("SELECT set_config('app.current_role', %s, true)", (role,))
                        
                        for idx, (chunk_text, embedding) in enumerate(
                            zip(page_chunks, embeddings)
                        ):
                            chunk_id = generate_chunk_id()
                            new_chunk_ids.append(chunk_id)
                            cur.execute(
                                """
                                INSERT INTO document_chunks
                                    (id, document_id, org_id, page_number,
                                     chunk_index, text, embedding,
                                     upload_mode, created_at)
                                VALUES (%s,%s,%s,%s,%s,%s,%s::vector,%s,%s)
                                """,
                                (
                                    chunk_id, registered_doc_id, org_id,
                                    page_num, idx, chunk_text,
                                    embedding, upload_mode.value,
                                    datetime.now(timezone.utc).isoformat()
                                )
                            )
                            chunks_created += 1
                    conn.commit()
                except Exception as e:
                    conn.rollback()
                    print(f"Neon chunk insert error: {e}")
                finally:
                    conn.close()

            store_page_hash(registered_doc_id, page_num, p_hash, new_chunk_ids, org_id )
            pages_newly_indexed += 1

        return {
            "doc_id": registered_doc_id,
            "is_new_document": False,
            "total_pages": total_pages,
            "pages_skipped": pages_skipped,
            "pages_newly_indexed": pages_newly_indexed,
            "chunks_created": chunks_created,
        }

    else:
        pages_newly_indexed = 0
        chunks_created = 0
        all_chunk_ids: List[str] = []

        conn = get_neon_connection()
        if conn:
            try:
                with conn.cursor() as cur:
                    # ---> ADDED RLS VARIABLES HERE <---
                    cur.execute("SELECT set_config('app.current_org_id', %s, true)", (org_id,))
                    cur.execute("SELECT set_config('app.current_user_id', %s, true)", (uploaded_by,))
                    cur.execute("SELECT set_config('app.current_role', %s, true)", (role,))
                    
                    for page_num, page_text in enumerate(pages, start=1):
                        p_hash = page_hash(page_text)
                        is_image = doc_type == DocumentType.IMAGE
                        page_chunks = [page_text] if is_image else chunk_by_type(doc_type, page_text)
                        embeddings = get_jina_embeddings(page_chunks, is_image=is_image)
                        page_chunk_ids: List[str] = []

                        for idx, (chunk_text, embedding) in enumerate(
                            zip(page_chunks, embeddings)
                        ):
                            chunk_id = generate_chunk_id()
                            page_chunk_ids.append(chunk_id)
                            all_chunk_ids.append(chunk_id)

                            cur.execute(
                                """
                                INSERT INTO document_chunks
                                    (id, document_id, org_id, page_number,
                                     chunk_index, text, embedding,
                                     upload_mode, created_at)
                                VALUES (%s,%s,%s,%s,%s,%s,%s::vector,%s,%s)
                                """,
                                (
                                    chunk_id, doc_id, org_id,
                                    page_num, idx, chunk_text,
                                    embedding, upload_mode.value,
                                    datetime.now(timezone.utc).isoformat()
                                )
                            )
                            chunks_created += 1
                        store_page_hash(doc_id, page_num, p_hash, page_chunk_ids, org_id)
                        pages_newly_indexed += 1

                conn.commit()
            except Exception as e:
                conn.rollback()
                print(f"Neon full ingestion error: {e}")
            finally:
                conn.close()

        store_file_in_registry(
            file_hash_val, org_id, doc_id,
            filename, total_pages,
            upload_mode.value, uploaded_by
        )

        return {
            "doc_id": doc_id,
            "is_new_document": True,
            "total_pages": total_pages,
            "pages_skipped": 0,
            "pages_newly_indexed": pages_newly_indexed,
            "chunks_created": chunks_created,
        }                                                                                                

def store_chunks_local_redis(
    doc_id: str, user_id: str, org_id: str, filename: str,
    chunks: List[str], embeddings: List[List[float]], doc_type: DocumentType
) -> int:
    if not redis_client:
        raise HTTPException(status_code=500, detail="Local session store (Redis) unavailable.")

    chunk_ids: List[str] = []
    pipeline = redis_client.pipeline()

    for idx, (chunk_text, embedding) in enumerate(zip(chunks, embeddings)):
        chunk_id = generate_chunk_id()
        chunk_ids.append(chunk_id)
        chunk_payload = json.dumps({
            "chunk_id": chunk_id, "doc_id": doc_id, "user_id": user_id,         
            "org_id": org_id, "file_name": filename, "doc_type": doc_type.value,
            "chunk_index": idx, "text": chunk_text, "embedding": embedding,
            "stored_at": datetime.now(timezone.utc).isoformat()
        })
        pipeline.setex(
            f"local:{user_id}:{doc_id}:chunk:{chunk_id}",
            LOCAL_SESSION_TTL, chunk_payload
        )

    pipeline.setex(
        f"local:{user_id}:{doc_id}:index",
        LOCAL_SESSION_TTL,
        json.dumps({
            "doc_id": doc_id, "file_name": filename, "chunk_ids": chunk_ids,
            "doc_type": doc_type.value, "created_at": datetime.now(timezone.utc).isoformat()
        })
    )

    user_docs_key = f"local:{user_id}:docs"
    existing_raw = redis_client.get(user_docs_key)
    existing_docs = json.loads(existing_raw) if existing_raw else []
    existing_docs.append(doc_id)

    pipeline.setex(user_docs_key, LOCAL_SESSION_TTL, json.dumps(existing_docs))
    pipeline.execute()
    return len(chunk_ids)

def log_local_upload_event_supabase(
    user_id: str, org_id: str, doc_id: str, filename: str, ip_address: str
) -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("audit_log").insert({
            "event_type": "local_upload_started",
            "user_id": user_id, "org_id": org_id, "doc_id": doc_id,
            "file_name": filename, "ip_address": ip_address,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "note": "Local session upload. Content stored ephemerally. No content logged."
        }).execute()
    except Exception as e:
        print(f"Audit log write failed: {e}")                                                                                                                   

global_upload_roles = RoleChecker([Role.ADMIN, Role.SUPER_ADMIN])
any_auth_user = RoleChecker([Role.USER, Role.ADMIN, Role.SUPER_ADMIN])                                                                                             

def _get_doc_name_from_registry(document_id: str) -> str:
    if not supabase_client:
        return document_id
    try:
        result = (
            supabase_client.table("document_registry")
            .select("file_name")
            .eq("id", document_id)
            .execute()
        )
        if result.data:
            return result.data[0]["file_name"]
    except Exception:
        pass
    return document_id

def retrieve_global_neon(
    query_embedding: List[float],
    org_id: str,
    user_id: str, 
    role: str, 
    top_k: int = 5,
    match_threshold: float = 0.2
) -> List[Dict[str, Any]]:
    conn = get_neon_connection()
    if not conn:
        print("Neon connection unavailable for global retrieval.")
        return []

    try:
        with conn.cursor(cursor_factory=RealDictCursor) as cur:
            
            cur.execute("SELECT set_config('app.current_org_id', %s, true)", (org_id,))
            cur.execute("SELECT set_config('app.current_user_id', %s, true)", (user_id,))
            cur.execute("SELECT set_config('app.current_role', %s, true)", (role,))
            
            cur.execute(
                """
                SELECT
                    id   AS chunk_id,
                    document_id,
                    org_id,
                    page_number,
                    chunk_index,
                    text,
                    upload_mode,
                    1 - (embedding <=> %s::vector) AS similarity_score
                FROM document_chunks
                WHERE org_id = %s
                  AND 1 - (embedding <=> %s::vector) >= %s
                ORDER BY similarity_score DESC
                LIMIT %s
                """,
                (query_embedding, org_id, query_embedding, match_threshold, top_k)
            )
            rows = cur.fetchall()

        results = []
        for row in rows:
            doc_name = _get_doc_name_from_registry(row["document_id"])
            results.append({
                "chunk_id":         row["chunk_id"],
                "document_id":      row["document_id"],
                "document_name":    doc_name,
                "org_id":           row["org_id"],
                "page_number":      row["page_number"],
                "chunk_index":      row["chunk_index"],
                "text":             row["text"],
                "similarity_score": float(row["similarity_score"]),
                "upload_mode":      row["upload_mode"],
            })
        return results

    except Exception as e:
        print(f"Neon retrieval error: {e}")
        return []
    finally:
        conn.close()

def retrieve_local_redis(
    query_embedding: List[float],
    user_id: str,
    org_id: str,
    top_k: int = 5,
    match_threshold: float = 0.1
) -> List[Dict[str, Any]]:
    if not redis_client:
        print("Redis unavailable for local retrieval.")
        return []
    try:
        chunk_keys = redis_client.keys(f"local:{user_id}:*:chunk:*")
        if not chunk_keys:
            return []
        scored: List[tuple] = []
        for key in chunk_keys:
            raw = redis_client.get(key)
            if not raw:
                continue
            item = json.loads(raw)
            if item.get("org_id") != org_id:
                continue
            stored_emb = item.get("embedding", [])
            if not stored_emb:
                continue
            score = sum(a * b for a, b in zip(query_embedding, stored_emb))
            if score >= match_threshold:
                scored.append((score, item))
        scored.sort(key=lambda x: x[0], reverse=True)
        top = scored[:top_k]
        results = []
        for score, item in top:
            results.append({
                "chunk_id": item.get("chunk_id", ""),
                "document_id": item.get("doc_id", ""),
                "document_name":item.get("file_name", "unknown"),
                "org_id": item.get("org_id", ""),
                "page_number": 0,                      
                "chunk_index": item.get("chunk_index", 0),
                "text": item.get("text", ""),
                "similarity_score": float(score),
                "upload_mode": "local",
            })
        return results
    except Exception as e:
        print(f"Redis retrieval error: {e}")
        return []
                                                                            
LANGUAGE_INSTRUCTION = {
    "English":    "Answer in English.", "Hindi":      "हिंदी में उत्तर दें।",
    "French":     "Répondez en français.", "German":     "Antworten Sie auf Deutsch.",
    "Spanish":    "Responde en español.", "Arabic":     "أجب باللغة العربية.",
    "Chinese":    "用中文回答。", "Japanese":   "日本語で答えてください。",
}

DEFAULT_SYSTEM_PROMPT = """You are an intelligent document assistant for an enterprise RAG platform.
Answer the user's question using ONLY the context provided below.
If the context does not contain enough information, say so clearly.
Do not fabricate or hallucinate any information not present in the context.
Be concise, accurate, and professional."""

def generate_llm_answer(
    user_query: str, context_chunks: List[Dict[str, Any]],
    language: str = "English", system_prompt: Optional[str] = None
) -> str:
    if not context_chunks:
        return (
            "I could not find relevant information in the available documents "
            "to answer your question. Please try rephrasing or upload relevant documents."
        )
    context_parts = []
    for i, chunk in enumerate(context_chunks, start=1):
        chunk_text = chunk.get('text', '')
        if isinstance(chunk_text, str) and chunk_text.startswith('data:image/'):
            context_parts.append(
                f"[Source {i} | Image: {chunk['document_name']}]\n(Image document matched via multimodal visual embedding: {chunk['document_name']})"
            )
        else:
            context_parts.append(
                f"[Source {i} | {chunk['document_name']} | Page {chunk['page_number']}]\n{chunk_text}"
            )
    context_text = "\n\n---\n\n".join(context_parts)
    lang_instruction = LANGUAGE_INSTRUCTION.get(language, "Answer in English.")
    base_prompt      = system_prompt or DEFAULT_SYSTEM_PROMPT
    full_system = f"{base_prompt}\n\nLanguage instruction: {lang_instruction}\n\nContext from documents:\n{context_text}"
    try:
        completion = groq_client.chat.completions.create(
            model = "qwen/qwen3.6-27b",
            messages=[
                {"role": "system", "content": full_system},
                {"role": "user",   "content": user_query}
            ],
            temperature=0.3,       
        )
        return completion.choices[0].message.content
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"LLM generation failed: {e}")


def save_chat_history(user_id: str, org_id: str, query: str, answer: str, query_mode: str, session_id:str) -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("chat_history").insert({
            "user_id": user_id,
            "session_id": session_id, 
            "org_id": org_id,
            "query": query,
            "answer": answer,
            "query_mode": query_mode,
            "created_at": datetime.now(timezone.utc).isoformat()
        }).execute()
    except Exception as e:
        print(f"Failed to save chat history: {e}")


def log_query_event(
    user_id: str, org_id: str, query: str,
    query_mode: str, sources_found: int, ip_address: str
) -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("audit_log").insert({
            "event_type":    "rag_query",
            "user_id":       user_id,
            "org_id":        org_id,
            "query_text":    query[:500],
            "query_mode":    query_mode,
            "sources_found": sources_found,
            "ip_address":    ip_address,
            "timestamp":     datetime.now(timezone.utc).isoformat()
        }).execute()
    except Exception as e:
        print(f"Query audit log failed: {e}")                                                                                      

@app.post("/api/v1/auth/signup", status_code=status.HTTP_201_CREATED)
async def signup(payload: UserSignup):
    if supabase_client:
        try:
            existing = (
                supabase_client.table("users")
                .select("id")
                .eq("email", payload.email)
                .execute()
            )
            if existing.data:
                raise HTTPException(status_code=400, detail="Email already registered")
        except HTTPException:
            raise
        except Exception as e:
            return {f"Supabase check failed: {e}"}

    user_id = f"usr_{uuid.uuid4().hex[:12]}"
    hashed_pw = hash_password(payload.password)
    org_id = payload.org_id or f"org_{uuid.uuid4().hex[:8]}"
    created_at = datetime.now(timezone.utc).isoformat()

    if supabase_client:
        try:
            supabase_client.table("users").insert({
                "id": user_id, "email": payload.email, "password_hash": hashed_pw,
                "role": payload.Role or Role.USER, "org_id": org_id,
                "tenant_id": payload.tenant_id, "created_at": created_at
            }).execute()
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"User creation failed: {e}")

    if redis_client:
        try:
            redis_client.setex(
                f"user:{user_id}:meta", 3600,
                json.dumps({"email": payload.email, "role": Role.USER, "org_id": org_id})
            )
        except Exception as e:
            return{f"Redis cache failed: {e}"}

    return {"user_id": user_id, "email": payload.email, "org_id": org_id} 

@app.post("/api/v1/auth/login", response_model=TokenResponse)
async def login(payload: UserLogin, response: Response):
    user = None
    if redis_client:
        try:
            cached = redis_client.get(f"user:email:{payload.email}")
            if cached:
                user = json.loads(cached)
        except Exception as e:
            print(f"Redis read failed: {e}")

    if not user and supabase_client:
        try:
            result = (
                supabase_client.table("users")
                .select("*")
                .eq("email", payload.email)
                .execute()
            )
            if not result.data:
                raise HTTPException(status_code=401, detail="Invalid credentials")
            user = result.data[0]

            if redis_client:
                redis_client.setex(f"user:email:{payload.email}", 3600, json.dumps(user))
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Auth lookup failed: {e}")

    if not user or not verify_password(payload.password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid credentials")

    token = create_jwt_token(
        claims={
            "user_id": user["id"], "email": user["email"], "role": user["role"],
            "org_id": user.get("org_id"), "tenant_id": user.get("tenant_id")
        },
        expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    )
    
    response.set_cookie(
        key="access_token", value=token, httponly=True,
        secure=True, samesite="lax", max_age=ACCESS_TOKEN_EXPIRE_MINUTES * 60
    )

    return TokenResponse(
        access_token=token, role=user["role"], org_id=user.get("org_id")
    )
                                        
@app.post("/api/v1/auth/logout")
async def logout(response: Response, current_user: TokenClaims = Depends(get_current_user)):
    if redis_client and current_user.exp:
        ttl = current_user.exp - int(datetime.now(timezone.utc).timestamp())
        if ttl > 0:
            redis_client.setex(f"blacklist:{current_user.user_id}", ttl, "1")
    response.delete_cookie("access_token")
    return {"message": "Logged out successfully"}

@app.get("/api/v1/auth/me", response_model=UserInfo)
async def get_me(current_user: TokenClaims = Depends(get_current_user)):
    return UserInfo(
        user_id=current_user.user_id, email=current_user.email, role=current_user.role,
        org_id=current_user.org_id, created_at=datetime.now(timezone.utc).isoformat()  
    )

@app.get("/api/v1/upload/consent", response_model=ConsentMessage)
async def get_upload_consent(
    upload_mode: UploadMode = Query(...), current_user: TokenClaims = Depends(get_current_user)
):
    if upload_mode == UploadMode.GLOBAL:
        if current_user.role not in [Role.ADMIN, Role.SUPER_ADMIN]:
            if not check_org_global_upload_setting(current_user.org_id):
                raise HTTPException(
                    status_code=status.HTTP_403_FORBIDDEN,
                    detail="Global uploads are disabled for standard users. Contact your Admin."
                )
        return ConsentMessage(
            upload_mode=UploadMode.GLOBAL, title="Global Upload Selected",
            message=("You have selected Global Upload. Your document will be permanently stored and made available to all authorised users within your organisation. It will remain accessible until manually removed by an Admin. Please ensure the document complies with your organisation's data retention and compliance policies."),
            confirm_label="Got it, Upload Globally", warning_label=None
        )
    return ConsentMessage(
        upload_mode=UploadMode.LOCAL, title="Local Upload Selected",
        message=("You have selected Local Upload. Your document will be stored privately for this session only — no other user, including Admins and Super Admins, can access or query it. After 1 hour (session expiry), your data will be automatically and permanently erased. It cannot be recovered after that point."),
        confirm_label="Got it, Upload Locally", warning_label="⚠️ Data will be erased after session ends (1 hour)."
    )                         

@app.post("/api/v1/upload/document", response_model=UploadReport, status_code=201)
async def upload_document(
    request: Request, file: UploadFile = File(...), upload_mode: UploadMode = Form(...),
    confirmed: bool = Form(...), current_user: TokenClaims = Depends(get_current_user)
):
    if not confirmed:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Upload requires explicit consent. Please confirm the upload mode disclaimer first."
        )

    if upload_mode == UploadMode.GLOBAL:
        if current_user.role not in [Role.ADMIN, Role.SUPER_ADMIN]:
            if not check_org_global_upload_setting(current_user.org_id):
                raise HTTPException(
                    status_code=status.HTTP_403_FORBIDDEN, 
                    detail="Global uploads are disabled for standard users. Contact your Admin."
                )

    org_id = current_user.org_id
    if not org_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN, detail="User has no organisation assigned. Contact your Admin."
        )

    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    file_hash_val = file_hash(file_bytes)
    doc_id        = generate_document_id()
    filename      = file.filename or "unknown"
    uploaded_at   = datetime.now(timezone.utc).isoformat()
    ip_address    = request.client.host if request.client else "unknown"

    doc_type = detect_document_type(filename, file_bytes)

    try:
        raw_text, total_pages = extract_text_by_type(doc_type, file_bytes, filename)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Content extraction failed for '{filename}': {e}")

    if doc_type == DocumentType.IMAGE:
        chunks = [raw_text]
    else:
        chunks = chunk_by_type(doc_type, raw_text)

    if not chunks:
        raise HTTPException(status_code=400, detail="No content could be extracted from the document.")

    if upload_mode == UploadMode.GLOBAL:

        store_raw_file_global_supabase(org_id, doc_id, filename, file_bytes)

    if upload_mode == UploadMode.GLOBAL:
        if supabase_client:
            try:
                supabase_client.table("audit_log").insert({
                    "event_type": "global_upload_started", "user_id": current_user.user_id,
                    "org_id": org_id, "doc_id": doc_id, "file_name": filename,
                    "file_hash": file_hash_val, "doc_type": doc_type.value,
                    "ip_address": ip_address, "timestamp": uploaded_at, "role": current_user.role
                }).execute()
            except Exception as e:
                print(f"Audit log error: {e}")

        
        delta = run_delta_management(
            doc_id=doc_id, org_id=org_id, filename=filename,
            file_hash_val=file_hash_val, raw_text=raw_text,
            doc_type=doc_type, chunks=chunks,
            uploaded_by=current_user.user_id, upload_mode=upload_mode,
            role=current_user.role.value 
        )

        return UploadReport(
            doc_id=delta["doc_id"], file_name=filename, upload_mode=upload_mode,
            org_id=org_id, doc_type=doc_type.value, total_pages=delta["total_pages"],
            pages_newly_indexed=delta["pages_newly_indexed"], pages_skipped=delta["pages_skipped"],
            chunks_created=delta["chunks_created"], status="indexed_globally" if delta["is_new_document"] else "delta_indexed",
            uploaded_at=uploaded_at
        )
    else:
        is_image = doc_type == DocumentType.IMAGE
        embeddings = get_jina_embeddings(chunks, is_image=is_image)
        store_raw_file_local_redis(current_user.user_id, doc_id, filename, file_bytes)
        chunks_created = store_chunks_local_redis(
            doc_id=doc_id, user_id=current_user.user_id, org_id=org_id, filename=filename,
            chunks=chunks, embeddings=embeddings, doc_type=doc_type
        )
        log_local_upload_event_supabase(
            user_id=current_user.user_id, org_id=org_id, doc_id=doc_id, filename=filename, ip_address=ip_address
        )
        return UploadReport(
            doc_id=doc_id, file_name=filename, upload_mode=upload_mode, org_id=org_id,
            doc_type=doc_type.value, total_pages=total_pages, pages_newly_indexed=total_pages, 
            pages_skipped=0, chunks_created=chunks_created, status="indexed_locally_session_ttl_1hr",
            uploaded_at=uploaded_at
        )                                                                                                                                                                                                                           


@app.get("/api/v1/documents/global/{doc_id}", tags=["Documents"])
async def get_global_document(doc_id: str, current_user: TokenClaims = Depends(get_current_user)):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client unavailable")
        
    response = supabase_client.table("document_registry").select("org_id, file_name").eq("id", doc_id).execute()
    if not response.data:
        raise HTTPException(status_code=404, detail="Document not found")
        
    org_id = response.data[0]["org_id"]
    filename = response.data[0]["file_name"]

    enforce_tenant_access(org_id, current_user)

    path = f"{org_id}/{doc_id}/{filename}"
    try:
        file_data = supabase_client.storage.from_("global_documents").download(path)
        media_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
        return Response(
            content=file_data, 
            media_type=media_type,
            headers={"Content-Disposition": f'inline; filename="{filename}"'}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error downloading from storage: {e}")

@app.get("/api/v1/documents/local/{doc_id}", tags=["Documents"])
async def get_local_document(doc_id: str, current_user: TokenClaims = Depends(get_current_user)):
    if not redis_client:
        raise HTTPException(status_code=500, detail="Redis client unavailable")
        
    raw = redis_client.get(f"local_raw:{current_user.user_id}:{doc_id}")
    if not raw:
        raise HTTPException(status_code=404, detail="Local document not found or session expired")
        
    data = json.loads(raw)
    file_bytes = base64.b64decode(data["data"])
    filename = data["filename"]
    media_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
    
    return Response(
        content=file_bytes, 
        media_type=media_type,
        headers={"Content-Disposition": f'inline; filename="{filename}"'}
    )


@app.post("/api/v1/query", response_model=RAGResponse)
async def query_rag(
    request: Request, body: RAGRequest, current_user: TokenClaims = Depends(get_current_user)
):
    org_id     = current_user.org_id
    user_id    = current_user.user_id
    ip_address = request.client.host if request.client else "unknown"

    if not org_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN, detail="User has no organisation assigned. Contact your Admin."
        )

    if not body.query or not body.query.strip():
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Query cannot be empty.")

    queried_at = datetime.now(timezone.utc).isoformat()

    try:
        query_embeddings = get_jina_embeddings([body.query], is_image=False)
        query_embedding  = query_embeddings[0]
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Query embedding failed: {e}")

    retrieved_chunks: List[Dict[str, Any]] = []

    if body.upload_mode == QueryMode.GLOBAL:
        
        retrieved_chunks = retrieve_global_neon(
            query_embedding=query_embedding, org_id=org_id, 
            user_id=user_id, role=current_user.role.value, top_k=body.top_k
        )

    elif body.upload_mode == QueryMode.LOCAL:
        retrieved_chunks = retrieve_local_redis(
            query_embedding=query_embedding, user_id=user_id, org_id=org_id, top_k=body.top_k
        )

    elif body.upload_mode == QueryMode.BOTH:
        local_chunks  = retrieve_local_redis(
            query_embedding=query_embedding, user_id=user_id, org_id=org_id, top_k=body.top_k
        )
        global_chunks = retrieve_global_neon(
            query_embedding=query_embedding, org_id=org_id,
            user_id=user_id, role=current_user.role.value, top_k=body.top_k
        )
        seen_ids: Dict[str, float] = {}
        merged: List[Dict] = []
        for chunk in local_chunks + global_chunks:
            cid   = chunk["chunk_id"]
            score = chunk["similarity_score"]
            if cid not in seen_ids or score > seen_ids[cid]:
                seen_ids[cid] = score
                merged.append(chunk)
        merged.sort(key=lambda x: x["similarity_score"], reverse=True)
        retrieved_chunks = merged[:body.top_k]

    current_session_id = body.session_id or f"session_{uuid.uuid4().hex[:12]}"

    answer = generate_llm_answer(
        user_query=body.query, context_chunks=retrieved_chunks,
        language=body.language, system_prompt=body.system_prompt
    )

    if body.upload_mode in [QueryMode.GLOBAL, QueryMode.BOTH]:
        save_chat_history(
            user_id=user_id,
            org_id=org_id,
            query=body.query,
            answer=answer,
            query_mode=body.upload_mode.value,
            session_id=current_session_id 
        )

    log_query_event(
        user_id=user_id, org_id=org_id, query=body.query,
        query_mode=body.upload_mode.value, sources_found=len(retrieved_chunks), ip_address=ip_address
    )

    sources = []
    base_url = str(request.base_url).rstrip("/")
    
    for c in retrieved_chunks:
        if c["upload_mode"] == "global":
            doc_url = f"{base_url}/api/v1/documents/global/{c['document_id']}"
        else:
            doc_url = f"{base_url}/api/v1/documents/local/{c['document_id']}"
            
        if c["document_name"].lower().endswith(".pdf"):
            doc_url += f"#page={c['page_number']}"

        raw_chunk_text = c.get("text", "")
        if isinstance(raw_chunk_text, str) and raw_chunk_text.startswith("data:image/"):
            preview_text = f"[Image Document: {c['document_name']}]"
        else:
            preview_text = raw_chunk_text[:200]

        sources.append(
            SourceResult(
                chunk_id=c["chunk_id"], 
                document_id=c["document_id"],
                document_name=c["document_name"],
                page_number=c["page_number"], 
                chunk_index=c["chunk_index"], 
                similarity_score=round(c["similarity_score"], 4),
                text_preview=preview_text, 
                org_id=c["org_id"], 
                upload_mode=c["upload_mode"],
                document_url=doc_url
            )
        )

    return RAGResponse(
        answer=answer, query=body.query, language=body.language,
        query_mode=body.upload_mode, sources=sources,
        total_sources_found=len(sources), generated_by="qwen/qwen3.6-27b (Groq)",
        org_id=org_id, queried_at=queried_at,
        session_id=current_session_id
    )

admin_only = RoleChecker([Role.ADMIN, Role.SUPER_ADMIN])
super_admin_only = RoleChecker([Role.SUPER_ADMIN]) 

@app.get("/api/v1/admin/users", tags=["Admin"])
async def get_organization_users(current_user: TokenClaims = Depends(admin_only)):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
    
    try:
        response = (
            supabase_client.table("users")
            .select("id, email, role, created_at")
            .eq("org_id", current_user.org_id)
            .execute()
        )
        return {"users": response.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch users: {e}")


@app.get("/api/v1/admin/settings/global-upload", response_model=OrgSettingsResponse, tags=["Admin"])
async def get_global_upload_setting(current_user: TokenClaims = Depends(any_auth_user)):

    is_allowed = check_org_global_upload_setting(current_user.org_id)
    return OrgSettingsResponse(
        org_id=current_user.org_id, 
        allow_user_global_upload=is_allowed
    )

@app.post("/api/v1/admin/settings/global-upload", tags=["Admin"])
async def configure_global_upload(
    payload: OrgSettingsUpdate,
    current_user: TokenClaims = Depends(admin_only)
):
    """Allows an Admin to toggle global upload access for standard users."""
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
    
    try:
        supabase_client.table("organization_settings").upsert({
            "org_id": current_user.org_id,
            "allow_user_global_upload": payload.allow_user_global_upload,
            "updated_at": datetime.now(timezone.utc).isoformat()
        }).execute()
        
        return {
            "message": "Global upload configuration updated successfully.", 
            "allow_user_global_upload": payload.allow_user_global_upload
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to update global upload settings: {e}")
    
@app.delete("/api/v1/admin/users/{target_user_id}", tags=["Admin"])
async def delete_organization_user(
    target_user_id: str, 
    current_user: TokenClaims = Depends(admin_only)
):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
        
    if current_user.user_id == target_user_id:
        raise HTTPException(status_code=400, detail="You cannot delete your own admin account.")

    try:
        user_check = supabase_client.table("users").select("org_id").eq("id", target_user_id).execute()
        if not user_check.data or user_check.data[0]["org_id"] != current_user.org_id:
            raise HTTPException(status_code=404, detail="User not found in your organization.")
            
        supabase_client.table("users").delete().eq("id", target_user_id).execute()
        
        if redis_client:
            redis_client.delete(f"user:{target_user_id}:meta")
            
        return {"message": f"User {target_user_id} removed successfully."}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to delete user: {e}")


@app.get("/api/v1/admin/documents", tags=["Admin"])
async def get_organization_documents(current_user: TokenClaims = Depends(admin_only)):
    """Allows an Admin to view all global documents in their organization."""
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
    
    try:
        response = (
            supabase_client.table("document_registry")
            .select("*")
            .eq("org_id", current_user.org_id)
            .execute()
        )
        return {"documents": response.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch documents: {e}")


@app.delete("/api/v1/admin/documents/{doc_id}", tags=["Admin"])
async def delete_organization_document(
    doc_id: str, 
    current_user: TokenClaims = Depends(admin_only)
):
    
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
        
    
    doc_check = supabase_client.table("document_registry").select("org_id").eq("id", doc_id).execute()
    if not doc_check.data or doc_check.data[0]["org_id"] != current_user.org_id:
        raise HTTPException(status_code=404, detail="Document not found in your organization.")

    
    conn = get_neon_connection()
    if conn:
        try:
            with conn.cursor() as cur:
                
                cur.execute(
                    "DELETE FROM document_chunks WHERE document_id = %s AND org_id = %s", 
                    (doc_id, current_user.org_id)
                )
                
                cur.execute(
                    "DELETE FROM image_store WHERE id = %s AND org_id = %s", 
                    (doc_id, current_user.org_id)
                )
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise HTTPException(status_code=500, detail=f"Failed to delete document chunks from DB: {e}")
        finally:
            conn.close()
    else:
        raise HTTPException(status_code=500, detail="Database connection failed.")

    
    try:
        
        supabase_client.table("page_registry").delete().eq("document_id", doc_id).execute()
        supabase_client.table("document_registry").delete().eq("id", doc_id).execute()
        
        
        supabase_client.table("audit_log").insert({
            "event_type": "global_document_deleted", 
            "user_id": current_user.user_id,
            "org_id": current_user.org_id, 
            "doc_id": doc_id,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }).execute()
        
        return {"message": f"Document {doc_id} and all related chunks successfully deleted."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to clear document registry: {e}")



@app.get("/api/v1/super-admin/users", tags=["Super Admin"])
async def super_admin_get_all_users(current_user: TokenClaims = Depends(super_admin_only)):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
    
    try:
        
        response = (
            supabase_client.table("users")
            .select("id, email, role, org_id, created_at")
            .execute()
        )
        return {"users": response.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch users globally: {e}")


@app.delete("/api/v1/super-admin/users/{target_user_id}", tags=["Super Admin"])
async def super_admin_delete_user(target_user_id: str, current_user: TokenClaims = Depends(super_admin_only)):
    
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
        
    if current_user.user_id == target_user_id:
        raise HTTPException(status_code=400, detail="You cannot delete your own Super Admin account.")

    try:
        
        supabase_client.table("users").delete().eq("id", target_user_id).execute()
        
        if redis_client:
            redis_client.delete(f"user:{target_user_id}:meta")
            
        return {"message": f"User {target_user_id} removed successfully from the platform."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to delete user globally: {e}")


@app.get("/api/v1/super-admin/documents", tags=["Super Admin"])
async def super_admin_get_all_documents(current_user: TokenClaims = Depends(super_admin_only)):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
    
    try:
        response = supabase_client.table("document_registry").select("*").execute()
        return {"documents": response.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch documents globally: {e}")


@app.delete("/api/v1/super-admin/documents/{doc_id}", tags=["Super Admin"])
async def super_admin_delete_document(doc_id: str, current_user: TokenClaims = Depends(super_admin_only)):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
        
    doc_check = supabase_client.table("document_registry").select("org_id").eq("id", doc_id).execute()
    if not doc_check.data:
        raise HTTPException(status_code=404, detail="Document not found on the platform.")
    
    target_org_id = doc_check.data[0]["org_id"]

    conn = get_neon_connection()
    if conn:
        try:
            with conn.cursor() as cur:
                cur.execute("DELETE FROM document_chunks WHERE document_id = %s", (doc_id,))
                cur.execute("DELETE FROM image_store WHERE id = %s", (doc_id,))
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise HTTPException(status_code=500, detail=f"Failed to delete document chunks globally from DB: {e}")
        finally:
            conn.close()
    else:
        raise HTTPException(status_code=500, detail="Database connection failed.")

    try:
        supabase_client.table("page_registry").delete().eq("document_id", doc_id).execute()
        supabase_client.table("document_registry").delete().eq("id", doc_id).execute()
        
        supabase_client.table("audit_log").insert({
            "event_type": "super_admin_global_document_deleted", 
            "user_id": current_user.user_id,
            "org_id": target_org_id, 
            "doc_id": doc_id,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }).execute()
        
        return {"message": f"Document {doc_id} and all related chunks successfully deleted globally."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to clear document registry globally: {e}")


@app.get("/api/v1/chat/history", tags=["Chat"])
async def get_chat_history(
    current_user: TokenClaims = Depends(get_current_user)
):
    if not supabase_client:
        raise HTTPException(status_code=500, detail="Supabase client not initialized")
    
    try:
        response = (
            supabase_client.table("chat_history")
            .select("*")
            .eq("user_id", current_user.user_id)
            .eq("org_id", current_user.org_id)
            .order("created_at", desc=True)
            .execute()
        )
        return {"history": response.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch chat history: {e}")

@app.get("/api/v1/guide", tags=["Help & Guide"])
async def platform_guide():

    contact_email = "mcode1929@gmail.com"

    return {
        "title": "Welcome to the Enterprise RAG Platform",
        "introduction": (
            "This platform acts as your intelligent document assistant. "
            "You can securely upload your documents (PDFs, Word docs, Spreadsheets, Presentations, etc.) "
            "and ask AI questions to instantly find answers based strictly on your files."
        ),
        "how_to_use_steps": [
            {
                "step": 1,
                "title": "🔐 Create an Account & Log In",
                "description": "Start by signing up with your email. Once logged in, you'll be securely assigned to your organization's workspace."
            },
            {
                "step": 2,
                "title": "📂 Upload Documents (Local vs. Global)",
                "description": "You have two ways to upload documents, depending on your needs:",
                "details": {
                    "Local Mode (Private & Temporary)": "Perfect for sensitive, one-off analysis. Documents are visible ONLY to you and are permanently deleted after 1 hour.",
                    "Global Mode (Org-Wide)": "Available for Admins. Documents uploaded globally act as a shared knowledge base for everyone in your organization."
                }
            },
            {
                "step": 3,
                "title": "💬 Ask Questions (Querying)",
                "description": "Head over to the chat interface to ask questions. You can filter where the AI searches for answers:",
                "details": {
                    "Local": "Searches only your temporarily uploaded files.",
                    "Global": "Searches your organization's permanent knowledge base.",
                    "Both": "Searches across both your private session files and the organization's files."
                }
            },
            {
                "step": 4,
                "title": "📜 View Chat History",
                "description": "Whenever you ask questions in 'Global' or 'Both' modes, your chat history is safely stored. You can revisit your past questions and answers at any time in the History tab."
            }
        ],
        "tips_for_best_results": [
            "Be specific with your questions.",
            "If the AI doesn't know the answer, it will tell you. It will never make up information.",
            "You can ask the AI to answer in different languages!"
        ],
        "support": "If you need elevated access (like Global Upload permissions), please contact your Organization's Admin.",
        "contact": f"if you see any issue or you have any suggetion then contact this email {contact_email} "
    }
    
@app.get("/health")
async def health():
    return {"status": "healthy"}

@app.get("/")
async def root():
    return {"service": "Global RAG v2.0", "services": ["local & global upload, RBAC , authentication , etc. "]}

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=7000)