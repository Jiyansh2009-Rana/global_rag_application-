import os
import io
import json
import base64
import mimetypes
import uuid
import logging
from datetime import datetime, timezone
from typing import List, Optional, Dict, Any, Tuple
import requests
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from pdf2image import convert_from_bytes
import pytesseract
from fastapi import HTTPException
from langchain_text_splitters import RecursiveCharacterTextSplitter

from core.config import JINA_API_KEY, LOCAL_SESSION_TTL, PAGES_PER_SET
from core.database import get_neon_connection, supabase_client, redis_client
from service.common.models import DocumentType, UploadMode, Role
from service.common.helper import (
    file_hash,
    page_hash,
    generate_chunk_id,
    generate_document_id,
    generate_set_id,
    extract_pages_as_list,
    split_pages_into_sets,
)

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENT TYPE DETECTION
# ─────────────────────────────────────────────────────────────────────────────
def detect_document_type(filename: str, file_bytes: bytes) -> DocumentType:
    filename_lower = filename.lower()
    if filename_lower.endswith('.txt') or filename_lower.endswith('.md'):
        return DocumentType.PLAIN_TEXT
    elif filename_lower.endswith('.pdf'):
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            first_page = reader.pages[0]
            text = first_page.extract_text()
            if not text or len(text.strip()) < 50:
                return DocumentType.PDF_SCANNED
            return DocumentType.PDF_TEXT
        except Exception:
            return DocumentType.PDF_SCANNED
    elif filename_lower.endswith('.docx'):
        return DocumentType.DOCX
    elif filename_lower.endswith(('.xlsx', '.csv')):
        return DocumentType.XLSX
    elif filename_lower.endswith('.pptx'):
        return DocumentType.PPTX
    elif filename_lower.endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp')):
        return DocumentType.IMAGE
    elif filename_lower.endswith('.html') or filename_lower.endswith('.htm'):
        return DocumentType.HTML
    return DocumentType.UNKNOWN

# ─────────────────────────────────────────────────────────────────────────────
# CHUNKING STRATEGIES
# ─────────────────────────────────────────────────────────────────────────────
def sentence_chunking(text: str) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n"]
    )
    return splitter.split_text(text)

def recursive_chunking(text: str) -> List[str]:
    recursive_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    return recursive_splitter.split_text(text)

def docx_chunking(text: str, max_chunk_size: int = 1500) -> List[str]:
    chunks = []
    current_section = ""
    for line in text.split('\n'):
        if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
            if current_section:
                chunks.append(current_section.strip())
            current_section = line
        else:
            current_section += "\n" + line
            if len(current_section) > max_chunk_size:
                chunks.append(current_section.strip())
                current_section = ""
    if current_section:
        chunks.append(current_section.strip())
    return [c for c in chunks if len(c) > 30]

def row_chunking(text: str, rows_per_chunk: int = 20) -> List[str]:
    if not text.strip():
        return []
    lines = text.split('\n')
    chunks = []
    current_chunk = []
    for line in lines:
        if not line.strip():
            continue
        current_chunk.append(line)
        if len(current_chunk) == rows_per_chunk:
            chunks.append("\n".join(current_chunk))
            current_chunk = []
    if current_chunk:
        chunks.append("\n".join(current_chunk))
    return chunks

def ppt_chunking(text: str) -> List[str]:
    slides = text.split('---SLIDE_BREAK---')
    return [s.strip() for s in slides if len(s.strip()) > 30]

def chunk_tag_aware(text: str, max_chunk_size: int = 1500) -> List[str]:
    chunks = []
    current_section = ""
    lines = [line for line in text.split('\n') if line.strip()]
    for line in lines:
        if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
            if current_section:
                chunks.append(current_section.strip())
            current_section = line
        else:
            current_section += "\n" + line
            if len(current_section) > max_chunk_size:
                chunks.append(current_section.strip())
                current_section = ""
    if current_section:
        chunks.append(current_section.strip())
    return [c for c in chunks if len(c) > 30]

# ─────────────────────────────────────────────────────────────────────────────
# FILE TEXT EXTRACTION
# ─────────────────────────────────────────────────────────────────────────────
def extract_pdf_text(file_bytes: bytes) -> Tuple[str, int]:
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        page_count = len(reader.pages)
        all_pages = []
        for page_number, page in enumerate(reader.pages):
            extracted_text = page.extract_text()
            if extracted_text and len(extracted_text.strip()) > 50:
                page_text = extracted_text.strip()
            else:
                try:
                    images = convert_from_bytes(
                        file_bytes,
                        dpi=300,
                        first_page=page_number + 1,
                        last_page=page_number + 1
                    )
                    if images:
                        page_text = pytesseract.image_to_string(
                            images[0],
                            lang="eng",
                            config="--oem 3 --psm 6"
                        ).strip()
                    else:
                        page_text = ""
                except Exception as ocr_err:
                    logger.warning(f"OCR failed for page {page_number + 1}: {ocr_err}")
                    page_text = extracted_text.strip() if extracted_text else ""
            all_pages.append(page_text)
        final_text = "\n---PAGE_BREAK---\n".join(all_pages)
        return final_text, page_count
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF extraction failed: {str(e)}")

def extract_docx_text(file_bytes: bytes) -> Tuple[str, int]:
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        text_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name.lower()
            if style_name.startswith('heading 1'):
                text = f"# {text}"
            elif style_name.startswith('heading 2'):
                text = f"## {text}"
            elif style_name.startswith('heading 3'):
                text = f"### {text}"
            text_parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                if any(row_data):
                    text_parts.append(" | ".join(row_data))
        full_text = "\n".join(text_parts)
        block_count = len(doc.paragraphs) + sum(len(t.rows) for t in doc.tables)
        return full_text, block_count
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX extraction failed: {str(e)}")

def extract_xlsx_text(file_bytes: bytes) -> Tuple[str, int]:
    full_text_lines = []
    total_rows = 0
    try:
        try:
            excel_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
        except Exception:
            excel_dict = {"Sheet1": pd.read_csv(io.BytesIO(file_bytes))}
        for sheet_name, df in excel_dict.items():
            df = df.fillna("")
            total_rows += len(df)
            columns = df.columns.tolist()
            for index, row in df.iterrows():
                row_parts = [f"Sheet: {sheet_name}"]
                for col in columns:
                    val = str(row[col]).strip()
                    if val:
                        row_parts.append(f"{col}: {val}")
                full_text_lines.append(" | ".join(row_parts))
        return "\n".join(full_text_lines), total_rows
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Spreadsheet extraction failed: {str(e)}")

def extract_text_from_shape(shape) -> str:
    text = ""
    if hasattr(shape, "text") and shape.text:
        text += shape.text + "\n"
    if shape.has_table:
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.replace('\n', ' ').strip()
                if cell_text:
                    row_data.append(cell_text)
            text += " | ".join(row_data) + "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            text += extract_text_from_shape(child_shape)
    return text

def extract_pptx_text(file_bytes: bytes) -> Tuple[str, int]:
    try:
        prs = PptxPresentation(io.BytesIO(file_bytes))
        full_text = ""
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                slide_text += extract_text_from_shape(shape)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text += f"\nSpeaker Notes:\n{notes}\n"
            full_text += slide_text + "\n---SLIDE_BREAK---\n"
        return full_text, len(prs.slides)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX extraction failed: {str(e)}")

def extract_html_text(file_bytes: bytes) -> Tuple[str, int]:
    try:
        html_text = file_bytes.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html_text, 'html.parser')
        for tag in soup.find_all(['nav', 'footer', 'script', 'style', 'aside', 'header', 'iframe']):
            tag.decompose()
        for level in range(1, 7):
            for header in soup.find_all(f'h{level}'):
                markdown_header = f"\n{'#' * level} {header.get_text(strip=True)}\n"
                header.string = markdown_header
        for table in soup.find_all('table'):
            table_lines = []
            for row in table.find_all('tr'):
                cells = row.find_all(['td', 'th'])
                row_data = [cell.get_text(strip=True).replace('\n', ' ') for cell in cells]
                if any(row_data):
                    table_lines.append(" | ".join(row_data))
            if table_lines:
                table.insert_after("\n" + "\n".join(table_lines) + "\n")
                table.decompose()
        text = soup.get_text(separator='\n', strip=True)
        return text, 1
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"HTML extraction failed: {str(e)}")

def extract_plainfile_text(file_bytes: bytes) -> Tuple[str, int]:
    text = file_bytes.decode('utf-8', errors='ignore')
    return text, 1

ALLOWED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp"}

IMAGE_EXTENSIONS = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}

def get_image_mime_type(filename: str) -> str:
    ext = os.path.splitext(filename.lower())[1]
    return IMAGE_EXTENSIONS.get(ext, mimetypes.guess_type(filename)[0] or "image/jpeg")

def read_image_to_base64(file_bytes: bytes, content_type: str) -> Tuple[str, int]:
    if content_type not in ALLOWED_IMAGE_TYPES:
        raise HTTPException(status_code=400, detail=f"Unsupported image type: {content_type}. Allowed: {ALLOWED_IMAGE_TYPES}")
    try:
        base64_encoded = base64.b64encode(file_bytes).decode("utf-8")
        data_url = f"data:{content_type};base64,{base64_encoded}"
        return data_url, 1
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Image encoding failed: {str(e)}")

def extract_text_by_type(doc_type: DocumentType, file_bytes: bytes, filename: str) -> Tuple[str, int]:
    if doc_type == DocumentType.PLAIN_TEXT:
        return extract_plainfile_text(file_bytes)
    elif doc_type in [DocumentType.PDF_TEXT, DocumentType.PDF_SCANNED]:
        return extract_pdf_text(file_bytes)
    elif doc_type == DocumentType.DOCX:
        return extract_docx_text(file_bytes)
    elif doc_type == DocumentType.XLSX:
        return extract_xlsx_text(file_bytes)
    elif doc_type == DocumentType.PPTX:
        return extract_pptx_text(file_bytes)
    elif doc_type == DocumentType.HTML:
        return extract_html_text(file_bytes)
    elif doc_type == DocumentType.IMAGE:
        mime_type = get_image_mime_type(filename)
        return read_image_to_base64(file_bytes, mime_type)
    else:
        text = file_bytes.decode('utf-8', errors='ignore')
        return text, 1

def chunk_by_type(doc_type: DocumentType, text: str) -> List[str]:
    if doc_type == DocumentType.IMAGE:
        return [text]
    elif doc_type == DocumentType.PLAIN_TEXT:
        return sentence_chunking(text)
    elif doc_type in [DocumentType.PDF_TEXT, DocumentType.PDF_SCANNED]:
        return recursive_chunking(text)
    elif doc_type == DocumentType.DOCX:
        return docx_chunking(text)
    elif doc_type == DocumentType.XLSX:
        return row_chunking(text, rows_per_chunk=20)
    elif doc_type == DocumentType.PPTX:
        return ppt_chunking(text)
    elif doc_type == DocumentType.HTML:
        return chunk_tag_aware(text)
    else:
        return recursive_chunking(text)

# ─────────────────────────────────────────────────────────────────────────────
# JINA EMBEDDINGS
# ─────────────────────────────────────────────────────────────────────────────
def get_jina_embeddings(texts: List[str], is_image: bool = False) -> List[List[float]]:
    url = "https://api.jina.ai/v1/embeddings"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {JINA_API_KEY}"
    }
    if is_image:
        input_data = [{"image": text} for text in texts]
    else:
        input_data = [{"text": text} for text in texts]
    payload = {
        "model": "jina-clip-v2",
        "normalized": True,
        "dimension": 1024,
        "input": input_data
    }
    try:
        response = requests.post(url, json=payload, headers=headers, timeout=30)
        if response.status_code == 200:
            res_json = response.json()
            return [item["embedding"] for item in res_json["data"]]
        else:
            raise HTTPException(status_code=502, detail=f"Jina embedding error: {response.text}")
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to generate embeddings from Jina API: {e}")

# ─────────────────────────────────────────────────────────────────────────────
# GLOBAL REGISTRY & DELTA ENGINE
# ─────────────────────────────────────────────────────────────────────────────
def check_file_in_registry(file_hash_val: str, org_id: str) -> Optional[Dict[str, Any]]:
    if not supabase_client:
        return None
    try:
        result = (
            supabase_client.table("document_registry")
            .select("*")
            .eq("file_hash", file_hash_val)
            .eq("org_id", org_id)
            .execute()
        )
        if result.data:
            return result.data[0]
    except Exception as e:
        logger.error(f"Registry lookup error: {e}")
    return None

def store_file_in_registry(
    file_hash_val: str, org_id: str, doc_id: str,
    file_name: str, total_pages: int,
    upload_mode: str, uploaded_by: str
) -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("document_registry").insert({
            "id": doc_id,
            "file_hash": file_hash_val,
            "org_id": org_id,
            "file_name": file_name,
            "total_pages": total_pages,
            "upload_mode": upload_mode,
            "uploaded_by": uploaded_by,
            "status": "ready",
            "uploaded_at": datetime.now(timezone.utc).isoformat()
        }).execute()
    except Exception as e:
        logger.error(f"Registry store error: {e}")

def check_page_hash_in_registry(doc_id: str, page_num: int, page_hash_val: str) -> bool:
    if not supabase_client:
        return False
    try:
        result = (
            supabase_client.table("page_registry")
            .select("id")
            .eq("document_id", doc_id)
            .eq("page_number", page_num)
            .eq("page_hash", page_hash_val)
            .execute()
        )
        return bool(result.data)
    except Exception:
        return False

def store_page_hash(
    doc_id: str, page_num: int, page_hash_val: str, chunk_ids: List[str], org_id: str = ""
) -> None:
    if not supabase_client:
        return
    try:
        supabase_client.table("page_registry").insert({
            "id": f"page_{uuid.uuid4().hex[:12]}",
            "document_id": doc_id,
            "page_number": page_num,
            "page_hash": page_hash_val,
            "chunk_ids": chunk_ids,
            "org_id": org_id,
            "indexed_at": datetime.now(timezone.utc).isoformat()
        }).execute()
    except Exception as e:
        logger.error(f"Page hash store error: {e}")

def store_raw_file_global_supabase(org_id: str, doc_id: str, filename: str, file_bytes: bytes) -> None:
    if not supabase_client:
        return
    try:
        path = f"{org_id}/{doc_id}/{filename}"
        supabase_client.storage.from_("global_documents").upload(
            file=file_bytes,
            path=path,
            file_options={"x-upsert": "true", "content-type": "application/octet-stream"}
        )
    except Exception as e:
        logger.error(f"Failed to store raw file in Supabase: {e}")

def process_set_global(
    set_id: str,
    doc_id: str,
    org_id: str,
    page_set: List[str],
    page_offsets: List[int],
    doc_type: DocumentType,
    uploaded_by: str,
    upload_mode: UploadMode,
    role: str,
) -> dict:
    is_image = doc_type == DocumentType.IMAGE
    all_chunk_records: List[tuple] = []
    page_hash_map: Dict[int, str] = {}

    for page_text, page_num in zip(page_set, page_offsets):
        p_hash = page_hash(page_text)
        if check_page_hash_in_registry(doc_id, page_num, p_hash):
            continue  # Page already indexed — skip

        if not page_text.strip():
            page_chunks = []
        else:
            page_chunks = [page_text] if is_image else chunk_by_type(doc_type, page_text)

        for idx, chunk_text in enumerate(page_chunks):
            all_chunk_records.append((chunk_text, page_num, idx))
        page_hash_map[page_num] = p_hash

    pages_newly_indexed = len(page_hash_map)
    pages_skipped = len(page_set) - pages_newly_indexed
    chunks_created = 0

    if not all_chunk_records:
        for page_num, p_hash_val in page_hash_map.items():
            store_page_hash(doc_id, page_num, p_hash_val, [], org_id)
        return {
            "set_id": set_id, "chunks_created": 0,
            "pages_newly_indexed": pages_newly_indexed, "pages_skipped": pages_skipped,
        }

    embeddings = get_jina_embeddings([r[0] for r in all_chunk_records], is_image=is_image)
    conn = get_neon_connection()
    if not conn:
        return {"set_id": set_id, "chunks_created": 0, "pages_newly_indexed": 0, "pages_skipped": pages_skipped}

    page_chunk_ids: Dict[int, List[str]] = {pn: [] for pn in page_hash_map}

    try:
        with conn.cursor() as cur:
            cur.execute("SELECT set_config('app.current_org_id', %s, true)", (org_id,))
            cur.execute("SELECT set_config('app.current_user_id', %s, true)", (uploaded_by,))
            cur.execute("SELECT set_config('app.current_role', %s, true)", (role,))

            for (chunk_text, page_num, chunk_idx), embedding in zip(all_chunk_records, embeddings):
                chunk_id = generate_chunk_id()
                page_chunk_ids[page_num].append(chunk_id)
                cur.execute(
                    """
                    INSERT INTO document_chunks
                        (id, document_id, org_id, page_number,
                         chunk_index, text, embedding,
                         upload_mode, set_id, created_at)
                    VALUES (%s,%s,%s,%s,%s,%s,%s::vector,%s,%s,%s)
                    """,
                    (
                        chunk_id, doc_id, org_id, page_num,
                        chunk_idx, chunk_text, embedding,
                        upload_mode.value, set_id,
                        datetime.now(timezone.utc).isoformat(),
                    ),
                )
                chunks_created += 1
        conn.commit()

        for page_num, p_hash_val in page_hash_map.items():
            store_page_hash(doc_id, page_num, p_hash_val, page_chunk_ids[page_num], org_id)

    except Exception as e:
        conn.rollback()
        logger.error(f"[set_global] {set_id} error: {e}")
    finally:
        conn.close()

    return {
        "set_id": set_id,
        "chunks_created": chunks_created,
        "pages_newly_indexed": pages_newly_indexed,
        "pages_skipped": pages_skipped,
    }

# ─────────────────────────────────────────────────────────────────────────────
# LOCAL REDIS DELTA & SESSION ENGINE
# ─────────────────────────────────────────────────────────────────────────────
def check_file_hash_local_redis(user_id: str, file_hash_val: str) -> Optional[str]:
    if not redis_client:
        return None
    try:
        raw = redis_client.get(f"local:filehash:{user_id}:{file_hash_val}")
        return raw
    except Exception:
        pass
    return None

def store_file_hash_local_redis(user_id: str, file_hash_val: str, doc_id: str, filename: str) -> None:
    if not redis_client:
        return
    try:
        redis_client.setex(
            f"local:filehash:{user_id}:{file_hash_val}",
            LOCAL_SESSION_TTL,
            json.dumps({"doc_id": doc_id, "file_name": filename}),
        )
    except Exception as e:
        logger.error(f"[local_delta] file hash store error: {e}")

def store_local_alias_audit_redis(
    user_id: str, org_id: str, doc_id: str,
    alias_filename: str,
    original_filename: str,
    file_hash: str,
) -> None:
    if not redis_client:
        return
    try:
        key = f"local:alias_audit:{user_id}:{doc_id}:{uuid.uuid4().hex[:8]}"
        redis_client.setex(key, LOCAL_SESSION_TTL, json.dumps({
            "event_type":        "duplicate_file_detected_local",
            "doc_id":            doc_id,
            "alias_filename":    alias_filename,
            "original_filename": original_filename,
            "file_hash":         file_hash,
            "user_id":           user_id,
            "org_id":            org_id,
            "timestamp":         datetime.now(timezone.utc).isoformat(),
        }))
    except Exception as e:
        logger.error(f"[local_alias_audit] Redis write failed: {e}")

def check_page_hash_local_redis(
    user_id: str, doc_id: str, page_num: int, page_hash_val: str
) -> bool:
    if not redis_client:
        return False
    try:
        stored = redis_client.get(f"local:delta:{user_id}:{doc_id}:page:{page_num}")
        return stored == page_hash_val
    except Exception:
        return False

def store_page_hash_local_redis(
    user_id: str, doc_id: str, page_num: int,
    page_hash_val: str, chunk_ids: List[str],
) -> None:
    if not redis_client:
        return
    try:
        pipe = redis_client.pipeline()
        pipe.setex(
            f"local:delta:{user_id}:{doc_id}:page:{page_num}",
            LOCAL_SESSION_TTL,
            page_hash_val,
        )
        pipe.setex(
            f"local:delta:{user_id}:{doc_id}:page:{page_num}:chunks",
            LOCAL_SESSION_TTL,
            json.dumps(chunk_ids),
        )
        pipe.execute()
    except Exception as e:
        logger.error(f"[local_delta] page hash store error: {e}")

def store_raw_file_local_redis(user_id: str, doc_id: str, filename: str, file_bytes: bytes) -> None:
    if not redis_client:
        return
    try:
        payload = json.dumps({
            "filename": filename,
            "data": base64.b64encode(file_bytes).decode("utf-8")
        })
        redis_client.setex(f"local_raw:{user_id}:{doc_id}", LOCAL_SESSION_TTL, payload)
    except Exception as e:
        logger.error(f"Failed to store raw file in Redis: {e}")

def process_set_local(
    set_id: str,
    doc_id: str,
    user_id: str,
    org_id: str,
    filename: str,
    page_set: List[str],
    page_offsets: List[int],
    doc_type: DocumentType,
) -> dict:
    if not redis_client:
        return {"set_id": set_id, "chunks_created": 0, "pages_newly_indexed": 0, "pages_skipped": len(page_set)}

    is_image = doc_type == DocumentType.IMAGE
    all_chunk_records: List[tuple] = []
    page_hash_map: Dict[int, str] = {}

    for page_text, page_num in zip(page_set, page_offsets):
        p_hash = page_hash(page_text)
        if check_page_hash_local_redis(user_id, doc_id, page_num, p_hash):
            continue

        if not page_text.strip():
            page_chunks = []
        else:
            page_chunks = [page_text] if is_image else chunk_by_type(doc_type, page_text)

        for idx, chunk_text in enumerate(page_chunks):
            all_chunk_records.append((chunk_text, page_num, idx))
        page_hash_map[page_num] = p_hash

    pages_newly_indexed = len(page_hash_map)
    pages_skipped = len(page_set) - pages_newly_indexed
    chunks_created = 0

    if not all_chunk_records:
        for page_num, p_hash_val in page_hash_map.items():
            store_page_hash_local_redis(user_id, doc_id, page_num, p_hash_val, [])
        return {
            "set_id": set_id, "chunks_created": 0,
            "pages_newly_indexed": pages_newly_indexed, "pages_skipped": pages_skipped,
        }

    embeddings = get_jina_embeddings([r[0] for r in all_chunk_records], is_image=is_image)
    page_chunk_ids: Dict[int, List[str]] = {pn: [] for pn in page_hash_map}

    try:
        pipeline = redis_client.pipeline()
        for (chunk_text, page_num, chunk_idx), embedding in zip(all_chunk_records, embeddings):
            chunk_id = generate_chunk_id()
            page_chunk_ids[page_num].append(chunk_id)
            chunk_payload = json.dumps({
                "chunk_id": chunk_id, "doc_id": doc_id, "user_id": user_id,
                "org_id": org_id, "file_name": filename, "doc_type": doc_type.value,
                "chunk_index": chunk_idx, "page_number": page_num, "text": chunk_text,
                "embedding": embedding, "set_id": set_id,
                "stored_at": datetime.now(timezone.utc).isoformat()
            })
            pipeline.setex(
                f"local:{user_id}:{doc_id}:chunk:{chunk_id}",
                LOCAL_SESSION_TTL, chunk_payload
            )
            chunks_created += 1

        pipeline.execute()

        for page_num, p_hash_val in page_hash_map.items():
            store_page_hash_local_redis(user_id, doc_id, page_num, p_hash_val, page_chunk_ids[page_num])

    except Exception as e:
        logger.error(f"[set_local] {set_id} error: {e}")

    return {
        "set_id": set_id,
        "chunks_created": chunks_created,
        "pages_newly_indexed": pages_newly_indexed,
        "pages_skipped": pages_skipped,
    }
